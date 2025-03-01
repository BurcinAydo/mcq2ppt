import re
from pptx import Presentation
from pptx.util import Pt, Inches

# Ask the user to paste text
print("Paste your multiple-choice test content below and press Enter twice to finish:")

input_text = []
while True:
    try:
        line = input()
        if line.strip() == "":
            break
        input_text.append(line)
    except EOFError:
        break

text = "\n".join(input_text)  # Preserve real new lines

# Ensure proper formatting of the text
text = text.replace("Question ", "\nQuestion ").replace("A) ", "\nA) ")

# Extract the test title
match = re.search(r"^(.*?)\nQuestion \d+:", text, re.DOTALL)
test_title = match.group(1).strip() if match else "Multiple-Choice Test"

# Extract questions and choices properly
questions = re.findall(r"(Question \d+:.*?)\n(A\).+?)(?=\nQuestion \d+:|\Z)", text, re.DOTALL)

# Create a new PowerPoint presentation
prs = Presentation()

# Add Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = test_title
subtitle.text = "Generated PowerPoint from Text Input"

# Format title slide font size
for shape in slide.shapes:
    if hasattr(shape, "text_frame") and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(28)

# Add each question as a separate slide
for question_text, choices_text in questions:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use a blank layout

    # Add Question
    question_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.5))
    question_frame = question_box.text_frame
    question_frame.word_wrap = True
    question_frame.auto_size = False
    question_frame.text = question_text.strip()
    question_frame.paragraphs[0].runs[0].font.size = Pt(30)  # Large question text

    # Add Answer Choices (Each on a New Line)
    choices_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
    choices_frame = choices_box.text_frame
    choices_frame.word_wrap = True
    choices_frame.auto_size = False
    choices_frame.clear()  # Ensure no default bullet formatting

    # **🚀 Final Fix: Forces Every a), b), c), d), e) onto a New Line**
    choices_text = re.sub(r"(?<!\n)([A-E])\)", lambda x: f"\n{x.group(1).lower()})", choices_text)  
    # Converts A) → a), B) → b), C) → c), etc.

    choices_list = choices_text.strip().split("\n")  # Split by real new lines

    # Add each choice separately
    for choice in choices_list:
        if choice.strip():
            p = choices_frame.add_paragraph()
            p.text = choice.strip()  # Ensure each choice starts on a new line
            p.space_after = Pt(10)  # Add spacing between choices
            run = p.add_run()  # Ensure at least one run exists
            run.font.size = Pt(24)  # Keep font at 24 pt for readability

# Save the presentation
pptx_filename = "multiple_choice_test.pptx"
prs.save(pptx_filename)

print(f"✅ PowerPoint file '{pptx_filename}' has been created successfully.")
